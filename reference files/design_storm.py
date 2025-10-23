#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import logging
import math
import re
from dataclasses import asdict, dataclass
from datetime import datetime, timedelta
from pathlib import Path
from typing import Iterable, Optional, Dict, Tuple

import numpy as np
import pandas as pd
import requests

try:
    import matplotlib.pyplot as plt  # optional for PNGs
except Exception:
    plt = None

# Optional NOAA helper lib
try:
    from pfdf.data.noaa import atlas14 as pfdf_atlas14  # type: ignore
except Exception:
    pfdf_atlas14 = None

try:
    import urllib.request
except Exception:
    urllib = None  # type: ignore


# =========================================================
# Temporal pattern sources
#   - PROPORTION_TABLES contains official NRCS/SCS dimensionless
#     cumulative rainfall curves for the SCS storm types.  These are
#     used directly for best accuracy.
#   - BETA_PRESETS provides Beta(α,β) shapes for the remaining patterns
#     (e.g., Huff quartiles).  SCS types remain in this mapping for
#     backward compatibility but are overridden when a dimensionless
#     table is available.
# =========================================================
BETA_PRESETS: Dict[str, Tuple[float, float]] = {
    # SCS/NRCS approximations (only used if table missing)
    "scs_type_i": (2.0, 5.0),
    "scs_type_ia": (2.0, 6.0),
    # Type II wants earlier peak; use skewed Beta (mode ≈ 0.33)
    # mode = (α−1)/(α+β−2) -> α=3.5, β=6 gives ~0.333
    "scs_type_ii": (3.5, 6.0),
    "scs_type_iii": (5.0, 2.0),

    # Huff quartiles (rough shapes; replace with your fitted values)
    "huff_q1": (1.5, 5.0),
    "huff_q2": (2.0, 3.0),
    "huff_q3": (3.0, 2.0),
    "huff_q4": (5.0, 1.5),

    # 'user' is handled separately via a custom cumulative CSV
    "user": (1.0, 1.0),
}



# Official NRCS/SCS dimensionless cumulative rainfall tables
PROPORTION_TABLES: Dict[str, np.ndarray] = {
    # Type I: 24-hour NRCS distribution (6-min increments).  Values increase
    # monotonically from 0 to 1.  Extracted from HydroCAD's Type I 24-hr
    # tabular file.  See design notes above for source.
    "scs_type_i": np.array([
        0.0000, 0.0017, 0.0035, 0.0052, 0.0070, 0.0087, 0.0105, 0.0122, 0.0139,
        0.0157, 0.0174, 0.0192, 0.0210, 0.0227, 0.0245, 0.0262, 0.0280, 0.0297,
        0.0315, 0.0332, 0.0350, 0.0368, 0.0386, 0.0404, 0.0423, 0.0442, 0.0461,
        0.0480, 0.0500, 0.0520, 0.0540, 0.0561, 0.0582, 0.0603, 0.0625, 0.0647,
        0.0669, 0.0691, 0.0714, 0.0737, 0.0760, 0.0784, 0.0807, 0.0831, 0.0855,
        0.0878, 0.0902, 0.0926, 0.0951, 0.0975, 0.1000, 0.1024, 0.1049, 0.1073,
        0.1098, 0.1123, 0.1148, 0.1174, 0.1199, 0.1225, 0.1250, 0.1276, 0.1303,
        0.1332, 0.1361, 0.1391, 0.1423, 0.1456, 0.1489, 0.1524, 0.1560, 0.1597,
        0.1633, 0.1671, 0.1708, 0.1746, 0.1784, 0.1823, 0.1861, 0.1901, 0.1940,
        0.1982, 0.2027, 0.2077, 0.2132, 0.2190, 0.2252, 0.2318, 0.2388, 0.2462,
        0.2540, 0.2623, 0.2714, 0.2812, 0.2917, 0.3030, 0.3194, 0.3454, 0.3878,
        0.4632, 0.5150, 0.5322, 0.5476, 0.5612, 0.5730, 0.5830, 0.5919, 0.6003,
        0.6083, 0.6159, 0.6230, 0.6298, 0.6365, 0.6430, 0.6493, 0.6555, 0.6615,
        0.6674, 0.6731, 0.6786, 0.6840, 0.6892, 0.6944, 0.6995, 0.7044, 0.7092,
        0.7140, 0.7186, 0.7232, 0.7276, 0.7320, 0.7362, 0.7404, 0.7444, 0.7484,
        0.7523, 0.7560, 0.7596, 0.7632, 0.7667, 0.7700, 0.7733, 0.7766, 0.7798,
        0.7830, 0.7862, 0.7894, 0.7926, 0.7958, 0.7989, 0.8020, 0.8051, 0.8082,
        0.8112, 0.8142, 0.8173, 0.8202, 0.8232, 0.8262, 0.8291, 0.8320, 0.8349,
        0.8378, 0.8406, 0.8434, 0.8462, 0.8490, 0.8518, 0.8546, 0.8573, 0.8600,
        0.8627, 0.8654, 0.8680, 0.8706, 0.8733, 0.8758, 0.8784, 0.8810, 0.8835,
        0.8860, 0.8885, 0.8910, 0.8934, 0.8958, 0.8982, 0.9006, 0.9030, 0.9054,
        0.9077, 0.9100, 0.9123, 0.9146, 0.9168, 0.9190, 0.9212, 0.9234, 0.9256,
        0.9278, 0.9299, 0.9320, 0.9341, 0.9362, 0.9382, 0.9402, 0.9423, 0.9442,
        0.9462, 0.9482, 0.9501, 0.9520, 0.9539, 0.9558, 0.9576, 0.9594, 0.9613,
        0.9630, 0.9648, 0.9666, 0.9683, 0.9700, 0.9717, 0.9734, 0.9750, 0.9766,
        0.9783, 0.9798, 0.9814, 0.9830, 0.9845, 0.9860, 0.9875, 0.9890, 0.9904,
        0.9918, 0.9933, 0.9946, 0.9960, 0.9974, 0.9987, 1.0000
    ]),
    # Type IA: 24-hour NRCS distribution in 0.5-hour increments (41 values)
    "scs_type_ia": np.array([
        0.0000, 0.0100, 0.0220, 0.0360, 0.0510,
        0.0670, 0.0830, 0.0990, 0.1160, 0.1350,
        0.1560, 0.1790, 0.2040, 0.2330, 0.2680,
        0.3100, 0.4250, 0.4800, 0.5200, 0.5500,
        0.5770, 0.6010, 0.6230, 0.6440, 0.6640,
        0.6830, 0.7010, 0.7190, 0.7360, 0.7530,
        0.7690, 0.7850, 0.8000, 0.8150, 0.8300,
        0.8440, 0.8580, 0.8710, 0.8840, 0.8960,
        0.9080, 0.9200, 0.9320, 0.9440, 0.9560,
        0.9670, 0.9780, 0.9890, 1.0000
    ]),
    # Type II: 24-hour NRCS distribution (6-min increments).  There are 241
    # values corresponding to 0.0, 0.1 hour, 0.2 hour, , 24 hours.  The final
    # value equals 1.0 exactly.  See HydroCAD table "Type II 24-hr Tabular".
    "scs_type_ii": np.array([
        0.0000, 0.0010, 0.0020, 0.0030, 0.0041, 0.0051, 0.0062, 0.0072, 0.0083, 0.0094,
        0.0105, 0.0116, 0.0127, 0.0138, 0.0150, 0.0161, 0.0173, 0.0184, 0.0196, 0.0208,
        0.0220, 0.0232, 0.0244, 0.0257, 0.0269, 0.0281, 0.0294, 0.0306, 0.0319, 0.0332,
        0.0345, 0.0358, 0.0371, 0.0384, 0.0398, 0.0411, 0.0425, 0.0439, 0.0452, 0.0466,
        0.0480, 0.0494, 0.0508, 0.0523, 0.0538, 0.0553, 0.0568, 0.0583, 0.0598, 0.0614,
        0.0630, 0.0646, 0.0662, 0.0679, 0.0696, 0.0712, 0.0730, 0.0747, 0.0764, 0.0782,
        0.0800, 0.0818, 0.0836, 0.0855, 0.0874, 0.0892, 0.0912, 0.0931, 0.0950, 0.0970,
        0.0990, 0.1010, 0.1030, 0.1051, 0.1072, 0.1093, 0.1114, 0.1135, 0.1156, 0.1178,
        0.1200, 0.1222, 0.1246, 0.1270, 0.1296, 0.1322, 0.1350, 0.1379, 0.1408, 0.1438,
        0.1470, 0.1502, 0.1534, 0.1566, 0.1598, 0.1630, 0.1663, 0.1697, 0.1733, 0.1771,
        0.1810, 0.1851, 0.1895, 0.1941, 0.1989, 0.2040, 0.2094, 0.2152, 0.2214, 0.2280,
        0.2350, 0.2427, 0.2513, 0.2609, 0.2715, 0.2830, 0.3068, 0.3544, 0.4308, 0.5679,
        0.6630, 0.6820, 0.6986, 0.7130, 0.7252, 0.7350, 0.7434, 0.7514, 0.7588, 0.7656,
        0.7720, 0.7780, 0.7836, 0.7890, 0.7942, 0.7990, 0.8036, 0.8080, 0.8122, 0.8162,
        0.8200, 0.8237, 0.8273, 0.8308, 0.8342, 0.8376, 0.8409, 0.8442, 0.8474, 0.8505,
        0.8535, 0.8565, 0.8594, 0.8622, 0.8649, 0.8676, 0.8702, 0.8728, 0.8753, 0.8777,
        0.8800, 0.8823, 0.8845, 0.8868, 0.8890, 0.8912, 0.8934, 0.8955, 0.8976, 0.8997,
        0.9018, 0.9038, 0.9058, 0.9078, 0.9097, 0.9117, 0.9136, 0.9155, 0.9173, 0.9192,
        0.9210, 0.9228, 0.9245, 0.9263, 0.9280, 0.9297, 0.9313, 0.9330, 0.9346, 0.9362,
        0.9377, 0.9393, 0.9408, 0.9423, 0.9438, 0.9452, 0.9466, 0.9480, 0.9493, 0.9507,
        0.9520, 0.9533, 0.9546, 0.9559, 0.9572, 0.9584, 0.9597, 0.9610, 0.9622, 0.9635,
        0.9647, 0.9660, 0.9672, 0.9685, 0.9697, 0.9709, 0.9722, 0.9734, 0.9746, 0.9758,
        0.9770, 0.9782, 0.9794, 0.9806, 0.9818, 0.9829, 0.9841, 0.9853, 0.9864, 0.9876,
        0.9887, 0.9899, 0.9910, 0.9922, 0.9933, 0.9944, 0.9956, 0.9967, 0.9978, 0.9989,
        1.0000
    ]),
    # Type III: 24-hour NRCS distribution (6-min increments).  Extracted from
    # HydroCAD's Type III 24-hr table.  Values range from 0 to 1.
    "scs_type_iii": np.array([
        0.0000, 0.0010, 0.0020, 0.0030, 0.0040, 0.0050, 0.0060, 0.0070, 0.0080, 0.0090,
        0.0100, 0.0110, 0.0120, 0.0130, 0.0140, 0.0150, 0.0160, 0.0170, 0.0180, 0.0190,
        0.0200, 0.0210, 0.0220, 0.0231, 0.0241, 0.0252, 0.0263, 0.0274, 0.0285, 0.0296,
        0.0308, 0.0319, 0.0331, 0.0343, 0.0355, 0.0367, 0.0379, 0.0392, 0.0404, 0.0417,
        0.0430, 0.0443, 0.0456, 0.0470, 0.0483, 0.0497, 0.0511, 0.0525, 0.0539, 0.0553,
        0.0567, 0.0582, 0.0597, 0.0612, 0.0627, 0.0642, 0.0657, 0.0673, 0.0688, 0.0704,
        0.0720, 0.0736, 0.0753, 0.0770, 0.0788, 0.0806, 0.0825, 0.0844, 0.0864, 0.0884,
        0.0905, 0.0926, 0.0948, 0.0970, 0.0993, 0.1016, 0.1040, 0.1064, 0.1089, 0.1114,
        0.1140, 0.1167, 0.1194, 0.1223, 0.1253, 0.1284, 0.1317, 0.1350, 0.1385, 0.1421,
        0.1458, 0.1496, 0.1535, 0.1575, 0.1617, 0.1659, 0.1703, 0.1748, 0.1794, 0.1842,
        0.1890, 0.1940, 0.1993, 0.2048, 0.2105, 0.2165, 0.2227, 0.2292, 0.2359, 0.2428,
        0.2500, 0.2578, 0.2664, 0.2760, 0.2866, 0.2980, 0.3143, 0.3394, 0.3733, 0.4166,
        0.5000, 0.5840, 0.6267, 0.6606, 0.6857, 0.7020, 0.7134, 0.7240, 0.7336, 0.7422,
        0.7500, 0.7572, 0.7641, 0.7708, 0.7773, 0.7835, 0.7895, 0.7952, 0.8007, 0.8060,
        0.8110, 0.8158, 0.8206, 0.8252, 0.8297, 0.8341, 0.8383, 0.8425, 0.8465, 0.8504,
        0.8543, 0.8579, 0.8615, 0.8650, 0.8683, 0.8716, 0.8747, 0.8777, 0.8806, 0.8833,
        0.8860, 0.8886, 0.8911, 0.8936, 0.8960, 0.8984, 0.9007, 0.9030, 0.9052, 0.9074,
        0.9095, 0.9116, 0.9136, 0.9156, 0.9175, 0.9194, 0.9212, 0.9230, 0.9247, 0.9264,
        0.9280, 0.9296, 0.9312, 0.9327, 0.9343, 0.9358, 0.9373, 0.9388, 0.9403, 0.9418,
        0.9433, 0.9447, 0.9461, 0.9475, 0.9489, 0.9503, 0.9517, 0.9530, 0.9544, 0.9557,
        0.9570, 0.9583, 0.9596, 0.9609, 0.9621, 0.9634, 0.9646, 0.9658, 0.9670, 0.9682,
        0.9694, 0.9706, 0.9718, 0.9729, 0.9741, 0.9752, 0.9764, 0.9775, 0.9786, 0.9797,
        0.9808, 0.9818, 0.9829, 0.9839, 0.9850, 0.9860, 0.9870, 0.9880, 0.9890, 0.9900,
        0.9909, 0.9919, 0.9928, 0.9938, 0.9947, 0.9956, 0.9965, 0.9974, 0.9983, 0.9991,
        1.0000
    ]),
}
# ---------------------------------------------------------------------
# (Optional) PROPORTION TABLE METHOD - commented out scaffolding.
def _storm_from_table(depth: float, n: int, table) -> np.ndarray:
    """Resample a dimensionless cumulative rainfall curve and return increments that sum to depth."""
    arr = np.asarray(list(table), dtype=float)
    if arr.ndim != 1 or arr.size < 2:
        raise ValueError("Dimensionless table must have at least two values")
    if not np.all(np.diff(arr) >= 0):
        raise ValueError("Dimensionless table must be non-decreasing")
    # Normalize to 0..1 and resample to n+1 points, then take diffs
    if arr[-1] == 0:
        arr = np.linspace(0, 1, arr.size)
    norm = arr / arr[-1]
    grid = np.linspace(0.0, 1.0, n + 1)
    cum = np.interp(grid, np.linspace(0.0, 1.0, norm.size), norm)
    inc = np.diff(cum)
    s = inc.sum()
    if not np.isfinite(s) or s <= 0:
        inc = np.full(n, depth / max(n, 1), dtype=float)
    else:
        inc = (inc / s) * depth
    return inc


def _beta_pdf(n: int, alpha: float, beta: float) -> np.ndarray:
    """Dimensionless pattern from Beta(,) on [0,1], no circular shift."""
    # mid-bin sampling to avoid spike at 0 and 1
    x = np.linspace(0, 1, n, endpoint=False) + 0.5 / n
    # Safe log evaluation for stability with fractional /
    with np.errstate(divide="ignore"):
        logpdf = (alpha - 1.0) * np.log(np.clip(x, 1e-12, 1.0)) + \
                 (beta  - 1.0) * np.log(np.clip(1.0 - x, 1e-12, 1.0))
    pdf = np.exp(logpdf)
    s = pdf.sum()
    if not np.isfinite(s) or s <= 0:
        pdf = np.full(n, 1.0 / n, dtype=float)
    else:
        pdf /= s
    return pdf


def beta_curve(n: int, alpha: float, beta: float, peak: int) -> np.ndarray:
    """Return a normalized Beta distribution with its peak at ``peak``.

    The underlying shape is defined by the ``alpha`` and ``beta`` parameters
    of the Beta distribution.  The resulting discrete curve is circularly
    shifted so that the maximum value occurs at index ``peak`` (0-based).
    This mirrors the behaviour of the historical ``beta_curve`` helper that
    some scripts depended on.
    """

    pdf = _beta_pdf(n, alpha, beta)
    if n <= 0:
        return pdf
    peak = int(peak)
    shift = (peak - int(np.argmax(pdf))) % n
    if shift:
        pdf = np.roll(pdf, shift)
    return pdf


def _user_pdf(n: int, path: Path) -> np.ndarray:
    """User CSV with [fraction/time, cumulative] columns (either can be dimensionless).
    Interprets first column as x in [0,1] or minutes normalized later; second as cumulative in [0,1].
    """
    df = pd.read_csv(path)
    t = np.asarray(df.iloc[:, 0], dtype=float)
    c = np.asarray(df.iloc[:, 1], dtype=float)
    # normalize if needed
    if t.max() > 1:
        t = t / t.max()
    if c.max() > 1:
        c = c / c.max()
    grid = np.linspace(0, 1, n + 1)
    cum = np.interp(grid, t, c)
    pdf = np.diff(cum)
    s = pdf.sum()
    if s <= 0 or not np.isfinite(s):
        pdf = np.full(n, 1.0 / n, dtype=float)
    else:
        pdf /= s
    return pdf


def build_storm(
    depth: float,
    duration_hr: float,
    timestep_min: float,
    distribution: str,
    peak: Optional[float] = None,
    custom_curve_path: Optional[Path] = None,
    start: Optional[datetime] = None,
) -> pd.DataFrame:
    """Generate a design storm DataFrame.

    Rules:
      - If ``distribution`` is one of the SCS/NRCS types present in PROPORTION_TABLES,
        resample that official dimensionless cumulative curve to the requested
        number of bins and scale by ``depth`` (best accuracy).
      - If ``distribution`` == "user", load a CSV with [time_fraction, cumulative_fraction]
        (either column may be dimensional and will be normalized to 0..1) and resample.
      - Otherwise, fall back to Beta(alpha, beta) presets from BETA_PRESETS (legacy).

    Output columns:
      time_min, intensity_in_hr, volume_in, cumulative_in (+ timestamp if ``start`` provided).
    """
    if duration_hr <= 0 or timestep_min <= 0:
        raise ValueError("Duration and timestep must be positive")

    # number of bins (ceil so last bin does not spill beyond duration)
    n = int(math.ceil((duration_hr * 60.0) / timestep_min))
    n = max(1, n)

    # choose temporal pattern
    if distribution in PROPORTION_TABLES:
        incremental = _storm_from_table(depth, n, PROPORTION_TABLES[distribution])
    elif distribution == "user":
        if not custom_curve_path:
            raise ValueError("Custom curve path required for 'user' distribution")
        pdf = _user_pdf(n, custom_curve_path)
        incremental = depth * pdf
    else:
        if distribution not in BETA_PRESETS:
            raise ValueError(f"Unknown distribution: {distribution}")
        a, b = BETA_PRESETS[distribution]
        pdf = _beta_pdf(n, a, b)
        incremental = depth * pdf

    # Intensities (in/hr); handle single-bin edge case
    if n == 1:
        intens = np.array([depth / max(duration_hr, 1e-12)], dtype=float)
    else:
        intens = incremental / (timestep_min / 60.0)

    cumulative = np.cumsum(incremental)
    minutes = np.arange(n) * timestep_min

    df = pd.DataFrame(
        {
            "time_min": minutes,
            "intensity_in_hr": intens,
            "volume_in": incremental,
            "cumulative_in": cumulative,
        }
    )
    if start is not None:
        df["timestamp"] = [start + timedelta(minutes=float(m)) for m in minutes]
    return df

# ------------------------ NOAA helpers -----------------------------------
_DURATION_RE = re.compile(
    r"""^\s*(?P<num>\d+(?:\.\d+)?)\s*[- ]\s*(?P<unit>min|minute|minutes|hr|hour|hours|day|days)\s*:?$""",
    re.IGNORECASE,
)


def _parse_noaa_text_to_df(txt: str) -> Optional[pd.DataFrame]:
    """Parse NOAA 'free text CSV' into a clean numeric DataFrame (rows=durations, cols=ARI years)."""
    lines = [ln.strip() for ln in txt.splitlines() if ln.strip()]
    header = next((ln for ln in lines if "ARI (years)" in ln), None)
    if not header:
        return None

    aris = re.findall(r"\b\d+\b", header.split("ARI (years)")[-1])
    cols = [str(int(x)) for x in aris if x.isdigit()]
    if not cols:
        return None

    data, idx = [], []
    for ln in lines:
        m = re.match(r"^([^:]+):\s*(.*)$", ln)
        if not m:
            continue
        label, rest = m.group(1).strip(), m.group(2)
        if not _DURATION_RE.match(label):
            continue  # skip Latitude:, Longitude:, etc.
        nums = re.findall(r"[-+]?(?:\d*\.\d+|\d+)(?:[eE][-+]?\d+)?", rest)
        nums = nums[: len(cols)] + [np.nan] * max(0, len(cols) - len(nums))
        row = []
        for x in nums:
            try:
                row.append(float(x))
            except Exception:
                row.append(np.nan)
        data.append(row)
        idx.append(label.rstrip(":"))

    if not data:
        return None
    df = pd.DataFrame(data, index=idx, columns=cols)
    return df


def _fetch_noaa_csv(stat: str, lat: float, lon: float) -> Optional[pd.DataFrame]:
    if urllib is None:
        return None
    try:
        base = f"https://hdsc.nws.noaa.gov/cgi-bin/new/fe_text_{stat}.csv"
        q = f"data=depth&lat={lat:.6f}&lon={lon:.6f}&series=pds&units=english"
        url = f"{base}?{q}"
        req = urllib.request.Request(url, headers={"User-Agent": "design-storm"})
        with urllib.request.urlopen(req, timeout=15) as r:
            txt = r.read().decode("utf-8", "replace")
        return _parse_noaa_text_to_df(txt)
    except Exception:
        return None


def fetch_noaa_table(lat: float, lon: float) -> Optional[pd.DataFrame]:
    # Prefer pfdf if installed; parse with same routine to keep behavior consistent
    if pfdf_atlas14 is not None:
        try:
            csv_path = pfdf_atlas14.download(
                lat, lon, statistic="mean", data="depth", series="pds", units="english", overwrite=True
            )
            txt = Path(csv_path).read_text(encoding="utf-8", errors="replace")
            df = _parse_noaa_text_to_df(txt)
            if df is not None and not df.empty:
                return df
        except requests.RequestException:
            logging.exception("NOAA table download failed")
            return None
    try:
        return _fetch_noaa_csv("mean", lat, lon)
    except requests.RequestException:
        logging.exception("NOAA table download failed")
        return None


def _label_to_minutes(label: str) -> float:
    m = _DURATION_RE.match(label.strip().rstrip(":"))
    if not m:
        return float("nan")
    num = float(m.group("num"))
    unit = m.group("unit").lower()
    if unit.startswith("min"):
        return num
    if unit.startswith("hr") or unit.startswith("hour"):
        return num * 60.0
    return num * 1440.0  # day


def _nearest_row_index(df: pd.DataFrame, target_minutes: float) -> int:
    mins = [_label_to_minutes(str(idx)) for idx in df.index]
    diffs = [abs(m - target_minutes) if pd.notna(m) else float("inf") for m in mins]
    return int(diffs.index(min(diffs)))


def fetch_noaa_depth(lat: float, lon: float, duration_hr: float, ari: float) -> Optional[float]:
    try_minutes = duration_hr * 60.0
    df = fetch_noaa_table(lat, lon)
    if df is None or df.empty:
        return None
    ari_key = str(int(round(ari)))
    if ari_key not in df.columns:
        return None
    row_i = _nearest_row_index(df, try_minutes)
    val = float(df.iloc[row_i][ari_key])
    return val if math.isfinite(val) else None


# ------------------------ presets/CLI ------------------------------------
@dataclass
class Config:
    location: Optional[str]
    duration: float
    return_period: float
    depth: Optional[float]
    time_step: float
    distribution: str
    peak: Optional[float]           # unused; kept for backward compatibility
    custom_curve: Optional[str]
    start_datetime: Optional[str]
    gauge_name: str
    export_type: str


def parse_args(argv: Iterable[str] | None = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Generate synthetic design storms using NRCS dimensionless tables or beta-distribution presets",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    p.add_argument("--location")
    p.add_argument("--duration", type=float, required=True)
    p.add_argument("--return-period", type=float, default=10.0)
    p.add_argument("--depth", type=float)
    p.add_argument("--time-step", type=float, default=5.0)
    p.add_argument(
        "--distribution",
        choices=list(BETA_PRESETS.keys()),
        default="scs_type_ii",
        help="Temporal pattern: SCS types use official NRCS dimensionless curves; others fall back to Beta(α,β) presets",
    )
    p.add_argument("--peak", type=float, default=None)  # ignored
    p.add_argument("--custom-curve", help="CSV of custom cumulative fractions (time, cum)")
    p.add_argument("--use-noaa", action="store_true")
    p.add_argument("--out-csv")
    p.add_argument("--out-dat")
    p.add_argument("--out-hyetograph")
    p.add_argument("--out-cumulative")
    p.add_argument("--pptx", default="", help="Path to PowerPoint file for generated plots")
    # SWMM export removed by request
    p.add_argument("--start-datetime")
    p.add_argument("--gauge-name", default="System")
    p.add_argument("--export-type", choices=["intensity", "volume", "cumulative"], default="intensity")
    p.add_argument("--save-preset")
    p.add_argument("--load-preset")
    p.add_argument(
        "-v", "--verbose", action="store_true", help="Increase logging verbosity"
    )
    p.add_argument(
        "-q", "--quiet", action="store_true", help="Suppress informational logs"
    )
    return p.parse_args(argv)


def save_preset(args, path):
    cfg = Config(
        args.location,
        args.duration,
        args.return_period,
        args.depth,
        args.time_step,
        args.distribution,
        args.peak,
        args.custom_curve,
        args.start_datetime,
        args.gauge_name,
        args.export_type,
    )
    Path(path).write_text(json.dumps(asdict(cfg), indent=2), encoding="utf-8")


def load_preset(path, args):
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    for k, v in data.items():
        if getattr(args, k, None) in (None,):
            setattr(args, k, v)


def plot_hyetograph(df: pd.DataFrame, path: Path) -> None:
    if plt is None:
        raise RuntimeError("matplotlib not available")
    fig, ax = plt.subplots()
    width = df["time_min"].diff().fillna(df["time_min"].iloc[1] if len(df) > 1 else 1.0)
    ax.bar(df["time_min"], df["intensity_in_hr"], width=width)
    ax.set_xlabel("Time (min)")
    ax.set_ylabel("Intensity (in/hr)")
    ax.set_title("Hyetograph")
    fig.tight_layout()
    fig.savefig(path)
    plt.close(fig)


def plot_cumulative(df: pd.DataFrame, path: Path) -> None:
    if plt is None:
        raise RuntimeError("matplotlib not available")
    fig, ax = plt.subplots()
    ax.plot(df["time_min"], df["cumulative_in"], marker="o")
    ax.set_xlabel("Time (min)")
    ax.set_ylabel("Cumulative Depth (in)")
    ax.set_title("Cumulative Mass Curve")
    fig.tight_layout()
    fig.savefig(path)
    plt.close(fig)


def add_image_to_pptx(pptx_path: Path, image_path: Path) -> None:
    """Append ``image_path`` as a slide in ``pptx_path``."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        logging.debug("pptx library not available; skipping PPTX export")
        return
    try:
        if pptx_path.exists():
            prs = Presentation(str(pptx_path))
        else:
            prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(0.5), width=Inches(9))
        prs.save(str(pptx_path))
        logging.info(f"Wrote PowerPoint -> {pptx_path}")
    except Exception as e:  # pragma: no cover - debug only
        logging.error(f"Failed to write PPTX: {e}")


def write_pcswmm_dat(
    df: pd.DataFrame,
    timestep_min: float,
    path: Path,
    gauge: str = "System",
    column: str = "intensity_in_hr",     # always in/hr per request
    start: Optional[datetime] = None,
) -> None:
    t0 = start or datetime(2003, 1, 1)
    with Path(path).open("w", encoding="utf-8", newline="") as f:
        f.write(";Rainfall (in/hr)\n;PCSWMM generated rain gauges file (please do not edit)\n")
        for i, val in enumerate(df[column], start=1):
            ts = t0 + timedelta(minutes=i * timestep_min)
            f.write(f"{gauge}\t{ts.year}\t{ts.month}\t{ts.day}\t{ts.hour}\t{ts.minute}\t{val:.7G}\n")


def main(argv: Iterable[str] | None = None) -> int:
    args = parse_args(argv)
    log_level = logging.INFO
    if args.quiet:
        log_level = logging.WARNING
    if args.verbose:
        log_level = logging.DEBUG
    logging.basicConfig(level=log_level, format="%(levelname)s: %(message)s")

    if args.use_noaa and args.location:
        try:
            lat, lon = map(float, args.location.split(","))
        except Exception:
            logging.error("Invalid --location 'lat,lon'")
            return 1
        depth = fetch_noaa_depth(lat, lon, args.duration, args.return_period)
        if depth is None and args.depth is None:
            logging.error("NOAA depth unavailable; provide --depth")
            return 1
        if depth is not None:
            args.depth = depth
            logging.info(f"NOAA depth: {depth:.3f} in")

    if args.depth is None:
        logging.error("Total depth is required")
        return 1

    custom_path = Path(args.custom_curve) if args.custom_curve else None
    start_dt = datetime.fromisoformat(args.start_datetime) if args.start_datetime else None

    df = build_storm(
        depth=args.depth,
        duration_hr=args.duration,
        timestep_min=args.time_step,
        distribution=args.distribution,
        peak=None,
        custom_curve_path=custom_path,
        start=start_dt,
    )

    col_map = {"intensity": "intensity_in_hr", "volume": "volume_in", "cumulative": "cumulative_in"}
    col = col_map[args.export_type]

    if args.out_csv:
        # Explicit line endings to prevent blank rows on Windows viewers
        Path(args.out_csv).with_suffix(".csv").write_text(
            df.to_csv(index=False, lineterminator="\r\n"), encoding="utf-8", newline=""
        )
    if args.out_dat:
        write_pcswmm_dat(df, args.time_step, Path(args.out_dat).with_suffix(".dat"), gauge=args.gauge_name, column="intensity_in_hr", start=start_dt)
    if args.out_hyetograph:
        plot_hyetograph(df, Path(args.out_hyetograph))
        if args.pptx:
            add_image_to_pptx(Path(args.pptx), Path(args.out_hyetograph))
    if args.out_cumulative:
        plot_cumulative(df, Path(args.out_cumulative))
        if args.pptx:
            add_image_to_pptx(Path(args.pptx), Path(args.out_cumulative))
    if args.save_preset:
        save_preset(args, args.save_preset)
    if not any([args.out_csv, args.out_dat, args.out_hyetograph, args.out_cumulative]):
        logging.info("\n%s", df.head())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())